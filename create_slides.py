"""
応用経済学会 プレゼンテーション スライド生成スクリプト
論文: International Standardization with Participation-Capacity Asymmetry
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour palette ──────────────────────────────────────────────
NAVY      = RGBColor(0x1B, 0x2A, 0x4A)
DARK_NAVY = RGBColor(0x0F, 0x1A, 0x30)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY= RGBColor(0xE8, 0xEB, 0xF0)
ACCENT    = RGBColor(0x3B, 0x82, 0xF6)   # bright blue accent
ACCENT2   = RGBColor(0x10, 0xB9, 0x81)   # teal-green accent
ACCENT3   = RGBColor(0xF5, 0x9E, 0x0B)   # amber accent
RED_ACC   = RGBColor(0xEF, 0x44, 0x44)
SUBTITLE_CLR = RGBColor(0x94, 0xA3, 0xB8)
BODY_BG   = RGBColor(0xF8, 0xFA, 0xFC)
CARD_BG   = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xE2, 0xE8, 0xF0)
CHARCOAL  = RGBColor(0x33, 0x41, 0x55)

FIGURES_DIR = os.path.join(os.path.dirname(__file__), "figures")

# ══════════════════════════════════════════════════════════════════
#  Helper Functions
# ══════════════════════════════════════════════════════════════════

def add_bg(slide, color=NAVY):
    """Fill slide background with a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    """Add a rounded-corner rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_accent_bar(slide, left, top, width=Inches(0.06), height=Inches(0.5), color=ACCENT):
    """Small vertical accent bar."""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=WHITE, alignment=PP_ALIGN.LEFT,
                 font_name="Meiryo", line_spacing=1.3):
    """Add a simple text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing:
        p.line_spacing = Pt(int(font_size * line_spacing))
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines, font_size=16,
                          color=WHITE, font_name="Meiryo", line_spacing=1.4,
                          alignment=PP_ALIGN.LEFT, bold=False):
    """Add text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(4)
        p.line_spacing = Pt(int(font_size * line_spacing))
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=15,
                    color=WHITE, bullet="●", font_name="Meiryo", line_spacing=1.5):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet}  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        p.line_spacing = Pt(int(font_size * line_spacing))
    return txBox


def add_slide_number(slide, number, total):
    """Small slide number at bottom-right."""
    add_text_box(slide, Inches(8.8), Inches(7.1), Inches(1.0), Inches(0.3),
                 f"{number} / {total}", font_size=10, color=SUBTITLE_CLR,
                 alignment=PP_ALIGN.RIGHT)


def add_section_header(slide, section_title, subtitle="", slide_num=0, total=0):
    """Standard section-header slide (dark background)."""
    add_bg(slide, DARK_NAVY)
    # Accent line
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(0.8), Inches(3.15), Inches(1.5), Inches(0.05))
    s.fill.solid()
    s.fill.fore_color.rgb = ACCENT
    s.line.fill.background()
    s.shadow.inherit = False

    add_text_box(slide, Inches(0.8), Inches(3.3), Inches(8.5), Inches(1.0),
                 section_title, font_size=32, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(4.2), Inches(8.5), Inches(0.6),
                     subtitle, font_size=16, color=SUBTITLE_CLR)
    if slide_num:
        add_slide_number(slide, slide_num, total)


def add_content_slide(slide, title, slide_num=0, total=0):
    """Set up a content slide with light background and title bar."""
    add_bg(slide, BODY_BG)
    # Top bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0), Inches(0), Inches(10), Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    bar.shadow.inherit = False

    add_text_box(slide, Inches(0.6), Inches(0.2), Inches(9.0), Inches(0.7),
                 title, font_size=22, bold=True, color=WHITE)
    if slide_num:
        add_slide_number(slide, slide_num, total)


# ══════════════════════════════════════════════════════════════════
#  Slide creation functions
# ══════════════════════════════════════════════════════════════════

TOTAL_SLIDES = 20

def slide_01_title(prs):
    """Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, DARK_NAVY)

    # Decorative top bar
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(0), Inches(0), Inches(10), Inches(0.08))
    s.fill.solid()
    s.fill.fore_color.rgb = ACCENT
    s.line.fill.background()
    s.shadow.inherit = False

    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(8.5), Inches(1.5),
                 "参加能力の非対称性を伴う\n国際標準化",
                 font_size=34, bold=True, color=WHITE, line_spacing=1.4)

    add_text_box(slide, Inches(0.8), Inches(3.3), Inches(8.5), Inches(0.6),
                 "International Standardization with Participation-Capacity Asymmetry",
                 font_size=15, color=SUBTITLE_CLR)

    # Separator
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(0.8), Inches(4.1), Inches(2.0), Inches(0.04))
    s.fill.solid()
    s.fill.fore_color.rgb = ACCENT
    s.line.fill.background()
    s.shadow.inherit = False

    add_text_box(slide, Inches(0.8), Inches(4.4), Inches(8.5), Inches(0.5),
                 "応用経済学会 報告", font_size=18, color=SUBTITLE_CLR)

    add_text_box(slide, Inches(0.8), Inches(5.1), Inches(8.5), Inches(0.5),
                 "松木 亮太", font_size=20, bold=True, color=WHITE)

    add_slide_number(slide, 1, TOTAL_SLIDES)


def slide_02_outline(prs):
    """Outline slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide(slide, "報告の構成", 2, TOTAL_SLIDES)

    items = [
        ("1.", "背景と動機", "なぜ「参加能力」に着目するのか"),
        ("2.", "モデルの設定", "3カ国・3レジーム・参加能力の非対称性"),
        ("3.", "企業均衡", "SW / SU / IS 下のクールノー均衡"),
        ("4.", "政府のレジーム選択", "均衡レジームと世界厚生の乖離"),
        ("5.", "能力投資と移転政策", "内生的参加能力とアクセッション条件"),
        ("6.", "まとめとインプリケーション", ""),
    ]

    y = Inches(1.4)
    for num, title, desc in items:
        # Number circle
        add_shape_rect(slide, Inches(0.7), y, Inches(0.45), Inches(0.45),
                       ACCENT, border_color=None)
        add_text_box(slide, Inches(0.7), y + Inches(0.02), Inches(0.45), Inches(0.4),
                     num, font_size=15, bold=True, color=WHITE,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(1.35), y, Inches(4.0), Inches(0.45),
                     title, font_size=17, bold=True, color=CHARCOAL)
        if desc:
            add_text_box(slide, Inches(1.35), y + Inches(0.35), Inches(7.5), Inches(0.3),
                         desc, font_size=12, color=SUBTITLE_CLR)
        y += Inches(0.85)


def slide_03_background(prs):
    """Background & motivation 1."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide(slide, "1. 背景：国際標準化の重要性", 3, TOTAL_SLIDES)

    items = [
        "技術標準は互換性・相互運用性・インストールドベースを通じ市場構造を規定",
        "ネットワーク外部性 → 標準化の成否が企業の革新・消費者の採用を左右\n（Farrell & Saloner 1985; Katz & Shapiro 1985–94）",
        "国際的文脈では、各国が自国標準の維持・外国標準の承認・共通標準化のいずれかを選択",
        "相互認証（MRA）・WTO/TBT・発展途上国の標準化参加など制度面の研究も蓄積\n（Gandal & Shy 2001; Barrett & Yang 2001; Costinot 2008; Klimenko 2009）",
    ]
    add_bullet_list(slide, Inches(0.7), Inches(1.3), Inches(8.6), Inches(5.5),
                    items, font_size=15, color=CHARCOAL, bullet="▸", line_spacing=1.6)


def slide_04_gap(prs):
    """Research gap slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide(slide, "1. 既存研究のギャップ", 4, TOTAL_SLIDES)

    # Left side - problem
    add_accent_bar(slide, Inches(0.7), Inches(1.4), height=Inches(1.5), color=RED_ACC)

    add_text_box(slide, Inches(0.95), Inches(1.4), Inches(4.0), Inches(0.4),
                 "既存理論の前提", font_size=16, bold=True, color=RED_ACC)
    add_multiline_textbox(slide, Inches(0.95), Inches(1.85), Inches(4.0), Inches(1.2),
                          ["共通標準への加盟（accession）＝",
                           "関連する摩擦の完全な除去"],
                          font_size=14, color=CHARCOAL)

    # Right side - reality
    add_accent_bar(slide, Inches(5.2), Inches(1.4), height=Inches(1.5), color=ACCENT2)

    add_text_box(slide, Inches(5.45), Inches(1.4), Inches(4.2), Inches(0.4),
                 "現実", font_size=16, bold=True, color=ACCENT2)
    add_multiline_textbox(slide, Inches(5.45), Inches(1.85), Inches(4.2), Inches(1.2),
                          ["形式的加盟 ≠ 実効的参加",
                           "検査・認証・適合性評価の能力は各国で非対称"],
                          font_size=14, color=CHARCOAL)

    # Arrow & key question
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(3.5),
                               Inches(8.6), Inches(0.04))
    s.fill.solid()
    s.fill.fore_color.rgb = ACCENT
    s.line.fill.background()
    s.shadow.inherit = False

    add_shape_rect(slide, Inches(0.7), Inches(3.8), Inches(8.6), Inches(1.6),
                   RGBColor(0xEF, 0xF6, 0xFF), border_color=ACCENT, border_width=Pt(1.5))
    add_text_box(slide, Inches(0.95), Inches(3.9), Inches(8.3), Inches(0.4),
                 "本論文の問い", font_size=16, bold=True, color=ACCENT)
    add_multiline_textbox(slide, Inches(0.95), Inches(4.3), Inches(8.3), Inches(1.0),
                          ["「参加能力の非対称性」がネットワーク外部性・連合形成と相互作用するとき、",
                           "国際標準化の均衡レジームと世界厚生はどう変わるか？"],
                          font_size=15, color=CHARCOAL)

    # Contribution summary at bottom
    add_text_box(slide, Inches(0.7), Inches(5.8), Inches(8.6), Inches(0.4),
                 "▸ 形式的互換性と実効的参加の区別を導入した初めてのモデル",
                 font_size=13, color=ACCENT, bold=True)

    add_slide_number(slide, 4, TOTAL_SLIDES)


def slide_05_model_overview(prs):
    """Model overview — section header."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "2. モデルの設定",
                       "3カ国・クールノー競争・参加能力の非対称性", 5, TOTAL_SLIDES)


def slide_06_three_regimes(prs):
    """Three regimes slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide(slide, "2. 三つの標準化レジーム", 6, TOTAL_SLIDES)

    regimes = [
        ("SW", "標準化戦争", "各国が自国標準を維持\n相互承認なし", RED_ACC,
         "外国市場コスト: c + τᵢ"),
        ("SU", "標準化連合", "A・B が相互承認\nC はブロック外", ACCENT3,
         "ブロック内: τᵢ\nブロック外: c + τᵢ"),
        ("IS", "国際標準化", "3カ国が共通標準を相互承認", ACCENT2,
         "全外国市場: τᵢ\n（c が消滅）"),
    ]

    x_start = Inches(0.5)
    for i, (abbr, title, desc, color, cost) in enumerate(regimes):
        x = x_start + Inches(i * 3.15)
        # Card
        add_shape_rect(slide, x, Inches(1.3), Inches(2.85), Inches(4.8),
                       CARD_BG, border_color=CARD_BORDER, border_width=Pt(1))

        # Colored header bar inside card
        hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      x, Inches(1.3), Inches(2.85), Inches(0.7))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = color
        hdr.line.fill.background()
        hdr.shadow.inherit = False

        add_text_box(slide, x + Inches(0.15), Inches(1.35), Inches(2.6), Inches(0.35),
                     abbr, font_size=20, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.15), Inches(1.65), Inches(2.6), Inches(0.3),
                     title, font_size=13, color=WHITE, alignment=PP_ALIGN.CENTER)

        # Description
        add_multiline_textbox(slide, x + Inches(0.2), Inches(2.2), Inches(2.5), Inches(1.5),
                              desc.split("\n"), font_size=13, color=CHARCOAL, line_spacing=1.5)

        # Cost box
        add_shape_rect(slide, x + Inches(0.15), Inches(3.9), Inches(2.55), Inches(1.5),
                       RGBColor(0xF1, 0xF5, 0xF9), border_color=CARD_BORDER)
        add_text_box(slide, x + Inches(0.25), Inches(3.95), Inches(2.4), Inches(0.3),
                     "外国市場限界費用", font_size=11, bold=True, color=ACCENT)
        add_multiline_textbox(slide, x + Inches(0.25), Inches(4.3), Inches(2.4), Inches(1.0),
                              cost.split("\n"), font_size=12, color=CHARCOAL)


def slide_07_key_equations(prs):
    """Key equations slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide(slide, "2. モデルの核心：参加能力と限界費用", 7, TOTAL_SLIDES)

    eqs = [
        ("残留コンプライアンス費用",
         "τₖ = μ(1 − θₖ)",
         "θₖ ∈ [0,1] は国 k の参加能力。高い θₖ ＝ 検査・認証が容易"),
        ("限界費用構造",
         "mᵢₖ = { 0  (自国),   τᵢ  (承認域内),   c + τᵢ  (承認域外) }",
         "相互承認は技術ギャップ c を除去するが、残留費用 τᵢ は残存"),
        ("逆需要関数",
         "pᵢₖ = 1 − Qₖ + v · B_G(i)",
         "ネットワーク外部性 v ≥ 0、互換的インストールドベース B_G(i)"),
        ("非対称性の仮定",
         "θ_A = θ_B > θ_C",
         "A・B は高能力、C は低能力 → τ_C が大きい"),
    ]

    y = Inches(1.25)
    for title, eq, note in eqs:
        # Card
        add_shape_rect(slide, Inches(0.6), y, Inches(8.8), Inches(1.15),
                       CARD_BG, border_color=CARD_BORDER, border_width=Pt(1))
        add_accent_bar(slide, Inches(0.6), y + Inches(0.1), height=Inches(0.95), color=ACCENT)

        add_text_box(slide, Inches(0.85), y + Inches(0.05), Inches(3.5), Inches(0.3),
                     title, font_size=13, bold=True, color=ACCENT)
        add_text_box(slide, Inches(0.85), y + Inches(0.35), Inches(8.3), Inches(0.35),
                     eq, font_size=15, bold=True, color=CHARCOAL, font_name="Consolas")
        add_text_box(slide, Inches(0.85), y + Inches(0.75), Inches(8.3), Inches(0.35),
                     note, font_size=12, color=SUBTITLE_CLR)
        y += Inches(1.3)


def slide_08_timing(prs):
    """Game timing slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide(slide, "2. ゲームのタイミング", 8, TOTAL_SLIDES)

    stages = [
        ("Stage 0", "政府が能力投資 Iₖ を選択", ACCENT2),
        ("Stage 1", "A・B が標準化連合を形成するか決定", ACCENT),
        ("Stage 2", "C がアクセッションを要請するか決定", ACCENT),
        ("Stage 3", "A・B が C の加盟を受諾するか決定", ACCENT),
        ("Stage 4", "企業がクールノー数量競争", ACCENT3),
    ]

    y = Inches(1.5)
    for label, desc, color in stages:
        # Stage badge
        add_shape_rect(slide, Inches(0.8), y, Inches(1.3), Inches(0.5),
                       color)
        add_text_box(slide, Inches(0.8), y + Inches(0.05), Inches(1.3), Inches(0.4),
                     label, font_size=14, bold=True, color=WHITE,
                     alignment=PP_ALIGN.CENTER)

        # Arrow connector
        arrow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(2.2), y + Inches(0.22), Inches(0.4), Inches(0.04))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = color
        arrow.line.fill.background()
        arrow.shadow.inherit = False

        # Desc
        add_text_box(slide, Inches(2.8), y + Inches(0.05), Inches(6.5), Inches(0.4),
                     desc, font_size=15, color=CHARCOAL)

        y += Inches(0.75)

    add_text_box(slide, Inches(0.8), y + Inches(0.4), Inches(8.4), Inches(0.4),
                 "▸ 後方帰納法（backward induction）により解を求める",
                 font_size=14, bold=True, color=ACCENT)


def slide_09_firm_eq_header(prs):
    """Firm equilibrium section header."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "3. 企業均衡",
                       "SW / SU / IS 各レジーム下のクールノー均衡", 9, TOTAL_SLIDES)


def slide_10_firm_eq_results(prs):
    """Key intuition of firm equilibrium."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide(slide, "3. 企業均衡の主要な特徴", 10, TOTAL_SLIDES)

    items = [
        "A・B の対称性を利用し、5つの均衡数量変数に集約：\n   x（自国生産）, y（相手メンバー国輸出）, z（C→A,B輸出）, w（A,B→C輸出）, t（C自国生産）",
        "SW: ネットワーク便益なし → 標準的なクールノー3社寡占\n   外国企業は c + τᵢ の限界費用ハンディ",
        "SU: A・B はインストールドベース S = 2X を通じてネットワーク便益を享受\n   C はブロック外のため便益なし → A・B の需要が外向きシフト",
        "IS: 全企業が共通インストールドベース N を共有\n   ただし τ_C が大きい → C の実効的参加は限定的",
    ]

    add_bullet_list(slide, Inches(0.6), Inches(1.2), Inches(8.8), Inches(5.5),
                    items, font_size=14, color=CHARCOAL, bullet="▸", line_spacing=1.7)


def slide_11_regime_schematic(prs):
    """Figure 1 - regime schematic."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide(slide, "3. レジーム間のコスト構造の比較", 11, TOTAL_SLIDES)

    fig_path = os.path.join(FIGURES_DIR, "fig1_regime_schematic.png")
    if os.path.exists(fig_path):
        slide.shapes.add_picture(fig_path, Inches(0.8), Inches(1.2), Inches(8.4))

    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(8.4), Inches(0.5),
                 "Figure 1: 各レジーム下の外国市場限界費用。SU/IS で c が除去されるが τᵢ は残存。",
                 font_size=11, color=SUBTITLE_CLR, alignment=PP_ALIGN.CENTER)
    add_slide_number(slide, 11, TOTAL_SLIDES)


def slide_12_govt_header(prs):
    """Government regime choice section header."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "4. 政府のレジーム選択",
                       "均衡レジーム・世界厚生との乖離・ネットワーク効果の非単調性",
                       12, TOTAL_SLIDES)


def slide_13_prop1(prs):
    """Proposition 1 - equilibrium regime."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide(slide, "4. 命題 1：均衡レジームの特徴付け", 13, TOTAL_SLIDES)

    # Main proposition box
    add_shape_rect(slide, Inches(0.5), Inches(1.2), Inches(9.0), Inches(3.3),
                   RGBColor(0xEF, 0xF6, 0xFF), border_color=ACCENT, border_width=Pt(2))

    add_text_box(slide, Inches(0.7), Inches(1.3), Inches(8.6), Inches(0.4),
                 "命題 1（Proposition 1）：均衡レジーム", font_size=16, bold=True, color=ACCENT)

    results = [
        "IS（国際標準化）⟺  Δ_C^I ≥ 0  かつ  Δ_A^I ≥ 0  かつ  Δ_A^{IS} ≥ 0",
        "SU（標準化連合）⟺  (Δ_C^I < 0  or  Δ_A^I < 0)  かつ  Δ_A^{SU} ≥ 0",
        "SW（標準化戦争）⟺  上記以外",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.85), Inches(8.4), Inches(2.5),
                    results, font_size=14, color=CHARCOAL, bullet="▸", line_spacing=1.7)

    # Key insight
    add_shape_rect(slide, Inches(0.5), Inches(4.8), Inches(9.0), Inches(1.6),
                   RGBColor(0xFE, 0xF3, 0xC7), border_color=ACCENT3, border_width=Pt(1.5))
    add_text_box(slide, Inches(0.7), Inches(4.9), Inches(8.6), Inches(0.35),
                 "核心的インサイト", font_size=14, bold=True, color=ACCENT3)
    add_multiline_textbox(slide, Inches(0.7), Inches(5.3), Inches(8.6), Inches(0.9),
                          ["Δ_A^I：メンバー国にとってISがSUより良いか（受諾条件）",
                           "Δ_C^I：アウトサイダーにとってISがSUより良いか（申請条件）",
                           "→ アクセッションには「C が望む」＋「A・B が受け入れる」の両方が必要"],
                          font_size=13, color=CHARCOAL)


def slide_14_prop2(prs):
    """Proposition 2 - welfare gap."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide(slide, "4. 命題 2：世界厚生と均衡レジームの乖離", 14, TOTAL_SLIDES)

    # Proposition box
    add_shape_rect(slide, Inches(0.5), Inches(1.2), Inches(9.0), Inches(2.0),
                   RGBColor(0xFE, 0xE2, 0xE2), border_color=RED_ACC, border_width=Pt(2))

    add_text_box(slide, Inches(0.7), Inches(1.3), Inches(8.6), Inches(0.4),
                 "命題 2：均衡レジームと世界最適の乖離", font_size=16, bold=True, color=RED_ACC)
    add_multiline_textbox(slide, Inches(0.7), Inches(1.8), Inches(8.6), Inches(1.2),
                          ["完全国際標準化（IS）が世界厚生を最大化するにもかかわらず、",
                           "均衡レジームが排他的2カ国ブロック（SU）に留まるパラメータ領域が存在する",
                           "検証点：(v,c,τ,τ_C) = (0.05, 0.05, 0.02, 0.18)"],
                          font_size=14, color=CHARCOAL)

    # Mechanism explanation
    add_text_box(slide, Inches(0.6), Inches(3.5), Inches(8.6), Inches(0.35),
                 "▸ メカニズム", font_size=15, bold=True, color=ACCENT)
    items = [
        "低参加能力の C が加盟しても、共通インストールドベースへの貢献は小さい",
        "一方、C の加盟はメンバー国市場で追加的競争圧力を生む",
        "→ メンバー国にとって IS は不利（Δ_A^I < 0）⇒ アクセッション拒否",
        "世界全体では IS の方が効率的でも、メンバー国の私的利益が排除を維持",
    ]
    add_bullet_list(slide, Inches(0.6), Inches(3.9), Inches(8.8), Inches(3.0),
                    items, font_size=13, color=CHARCOAL, bullet="●", line_spacing=1.6)


def slide_15_fig_divergence(prs):
    """Figure 3 - divergence map."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide(slide, "4. 厚生–均衡乖離の視覚化", 15, TOTAL_SLIDES)

    fig_path = os.path.join(FIGURES_DIR, "fig3_divergence_map.png")
    if os.path.exists(fig_path):
        slide.shapes.add_picture(fig_path, Inches(1.0), Inches(1.2), Inches(7.5))

    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(8.4), Inches(0.5),
                 "Figure 3: 均衡がSUだが世界最適はISであるパラメータ領域（シェード部分）",
                 font_size=11, color=SUBTITLE_CLR, alignment=PP_ALIGN.CENTER)
    add_slide_number(slide, 15, TOTAL_SLIDES)


def slide_16_prop3(prs):
    """Proposition 3 - network effects."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide(slide, "4. 命題 3：ネットワーク効果の非単調性", 16, TOTAL_SLIDES)

    # Main box
    add_shape_rect(slide, Inches(0.5), Inches(1.2), Inches(9.0), Inches(1.7),
                   RGBColor(0xD1, 0xFA, 0xE5), border_color=ACCENT2, border_width=Pt(2))
    add_text_box(slide, Inches(0.7), Inches(1.3), Inches(8.6), Inches(0.35),
                 "命題 3：より強いネットワーク効果が排他的ブロックを安定化させうる",
                 font_size=15, bold=True, color=ACCENT2)
    add_multiline_textbox(slide, Inches(0.7), Inches(1.75), Inches(8.6), Inches(0.9),
                          ["v' < v'' なる v 値が存在し、v' では IS が均衡だが v'' では SU が均衡",
                           "検証例：(c,τ,τ_C) = (0.05, 0, 0.15)、v'=0.04 → IS、v''=0.055 → SU"],
                          font_size=13, color=CHARCOAL)

    # Intuition
    add_text_box(slide, Inches(0.6), Inches(3.2), Inches(8.6), Inches(0.35),
                 "▸ 直感的説明", font_size=15, bold=True, color=ACCENT)
    items = [
        "ネットワーク効果 v の増大 → 共通承認のインストールドベース拡大の利得増",
        "同時に、既存ブロック（A・B）の戦略的価値も増大",
        "τ_C が高い場合、C の効果的参加が弱いため、ブロック拡大の利得 < ブロック維持の利得",
        "→ v の増大がむしろ排他的ブロックを安定化させるパラドックス",
    ]
    add_bullet_list(slide, Inches(0.6), Inches(3.6), Inches(8.8), Inches(3.0),
                    items, font_size=13, color=CHARCOAL, bullet="●", line_spacing=1.6)


def slide_17_investment_header(prs):
    """Section 5 header."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "5. 能力投資と移転政策",
                       "参加能力の内生化・アクセッション閾値・条件付き移転",
                       17, TOTAL_SLIDES)


def slide_18_investment(prs):
    """Capacity investment and transfer policy."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide(slide, "5. 能力投資と条件付き移転", 18, TOTAL_SLIDES)

    # Investment box
    add_shape_rect(slide, Inches(0.5), Inches(1.2), Inches(4.3), Inches(2.6),
                   CARD_BG, border_color=ACCENT, border_width=Pt(1.5))
    add_text_box(slide, Inches(0.65), Inches(1.3), Inches(4.0), Inches(0.35),
                 "能力投資（Stage 0）", font_size=15, bold=True, color=ACCENT)
    items_inv = [
        "θₖ = θ̄ₖ + Iₖ　（投資で能力向上）",
        "費用：C(Iₖ) = (η/2)Iₖ²",
        "アウトサイダー投資 → アクセッション制約の緩和",
        "メンバー投資 → ブロック形成制約の緩和",
    ]
    add_bullet_list(slide, Inches(0.65), Inches(1.75), Inches(4.0), Inches(2.0),
                    items_inv, font_size=12, color=CHARCOAL, bullet="▸", line_spacing=1.6)

    # Transfer box
    add_shape_rect(slide, Inches(5.2), Inches(1.2), Inches(4.3), Inches(2.6),
                   CARD_BG, border_color=ACCENT2, border_width=Pt(1.5))
    add_text_box(slide, Inches(5.35), Inches(1.3), Inches(4.0), Inches(0.35),
                 "条件付き移転（Conditional Transfer）", font_size=15, bold=True, color=ACCENT2)
    items_tr = [
        "IS 実現時のみ有効な移転 s ≥ 0",
        "C の実効コスト：τ_C^IS(s) = τ_C − s",
        "移転費用：G(s) = (κ/2)s²",
        "プランナーが最適 s を選択",
    ]
    add_bullet_list(slide, Inches(5.35), Inches(1.75), Inches(4.0), Inches(2.0),
                    items_tr, font_size=12, color=CHARCOAL, bullet="▸", line_spacing=1.6)

    # Key result
    add_shape_rect(slide, Inches(0.5), Inches(4.1), Inches(9.0), Inches(2.3),
                   RGBColor(0xEF, 0xF6, 0xFF), border_color=ACCENT, border_width=Pt(1.5))
    add_text_box(slide, Inches(0.7), Inches(4.2), Inches(8.6), Inches(0.35),
                 "主要結果（命題 4）", font_size=15, bold=True, color=ACCENT)
    items_res = [
        "最小アクセッション移転 s^acc = max{0, τ_C − τ_C^I(τ)} が鍵",
        "W^IS(τ, τ_C − s^acc) − G(s^acc) ≥ max{W^SU, W^SW} ならば IS を誘導可能",
        "断片化が参加能力の弱さに起因する場合、能力構築型の移転が\n標準フラグメンテーション維持より効果的",
    ]
    add_bullet_list(slide, Inches(0.7), Inches(4.65), Inches(8.6), Inches(1.6),
                    items_res, font_size=13, color=CHARCOAL, bullet="●", line_spacing=1.6)


def slide_19_fig_investment(prs):
    """Figure 4 - investment map."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide(slide, "5. 政策空間でのレジームマップ", 19, TOTAL_SLIDES)

    fig_path = os.path.join(FIGURES_DIR, "fig4_investment_map.png")
    if os.path.exists(fig_path):
        slide.shapes.add_picture(fig_path, Inches(1.0), Inches(1.2), Inches(7.5))

    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(8.4), Inches(0.5),
                 "Figure 4: (I_H, I_C) 空間でのレジーム境界。アウトサイダー投資がアクセッション制約を緩和。",
                 font_size=11, color=SUBTITLE_CLR, alignment=PP_ALIGN.CENTER)
    add_slide_number(slide, 19, TOTAL_SLIDES)


def slide_20_conclusion(prs):
    """Conclusion slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_NAVY)

    # Accent bar
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(0), Inches(0), Inches(10), Inches(0.08))
    s.fill.solid()
    s.fill.fore_color.rgb = ACCENT
    s.line.fill.background()
    s.shadow.inherit = False

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8.5), Inches(0.6),
                 "6. まとめとインプリケーション", font_size=26, bold=True, color=WHITE)

    # Three main results
    results = [
        ("結果 1", "IS が世界最適でも均衡は SU になりうる",
         "低参加能力のアウトサイダーはインストールドベースへの\n貢献が小さく、メンバー国がアクセッションを拒否", RED_ACC),
        ("結果 2", "ネットワーク効果は排他的ブロックを安定化させうる",
         "参加能力の非対称性が十分大きいとき、v の増大は\nむしろ2カ国ブロックを強化", ACCENT3),
        ("結果 3", "能力構築型政策が最も効果的",
         "標準フラグメンテーション維持より、検査・認証・\n適合性評価への投資が真のボトルネックを解消", ACCENT2),
    ]

    y = Inches(1.4)
    for label, title, desc, color in results:
        # Label badge
        add_shape_rect(slide, Inches(0.8), y, Inches(1.1), Inches(0.4), color)
        add_text_box(slide, Inches(0.8), y + Inches(0.02), Inches(1.1), Inches(0.35),
                     label, font_size=13, bold=True, color=WHITE,
                     alignment=PP_ALIGN.CENTER)

        add_text_box(slide, Inches(2.1), y, Inches(7.3), Inches(0.35),
                     title, font_size=15, bold=True, color=WHITE)

        add_multiline_textbox(slide, Inches(2.1), y + Inches(0.4), Inches(7.3), Inches(0.8),
                              desc.split("\n"), font_size=12, color=SUBTITLE_CLR)
        y += Inches(1.45)

    # Policy message
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.9),
                               Inches(8.4), Inches(0.04))
    s.fill.solid()
    s.fill.fore_color.rgb = ACCENT
    s.line.fill.background()
    s.shadow.inherit = False

    add_text_box(slide, Inches(0.8), Inches(6.1), Inches(8.4), Inches(0.35),
                 "政策含意", font_size=16, bold=True, color=ACCENT)
    add_multiline_textbox(slide, Inches(0.8), Inches(6.5), Inches(8.4), Inches(0.8),
                          ["断片化の原因が「参加能力の弱さ」にあるなら、",
                           "形式的承認だけでなく、能力構築（testing, certification, conformity assessment）が鍵"],
                          font_size=13, color=WHITE)

    add_slide_number(slide, 20, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════════
#  Main
# ══════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_01_title(prs)
    slide_02_outline(prs)
    slide_03_background(prs)
    slide_04_gap(prs)
    slide_05_model_overview(prs)
    slide_06_three_regimes(prs)
    slide_07_key_equations(prs)
    slide_08_timing(prs)
    slide_09_firm_eq_header(prs)
    slide_10_firm_eq_results(prs)
    slide_11_regime_schematic(prs)
    slide_12_govt_header(prs)
    slide_13_prop1(prs)
    slide_14_prop2(prs)
    slide_15_fig_divergence(prs)
    slide_16_prop3(prs)
    slide_17_investment_header(prs)
    slide_18_investment(prs)
    slide_19_fig_investment(prs)
    slide_20_conclusion(prs)

    out_path = os.path.join(os.path.dirname(__file__), "presentation_slides.pptx")
    prs.save(out_path)
    print(f"Saved: {out_path}")


if __name__ == "__main__":
    main()
